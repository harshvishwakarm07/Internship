from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

OUT_PATH = Path("d:/InternShip/docs/SITS_Teacher_Submission.pptx")

slides_data = [
    ("Student Internship Tracking System (SITS)", [
        "Full-Stack MERN Project",
        "Prepared for Teacher Evaluation",
        "Prepared by: ____________________",
        "Submission Date: ____________________",
    ]),
    ("Agenda", [
        "Problem Statement",
        "Objectives",
        "Architecture",
        "DFD and ER Diagram",
        "Database Models and Schema",
        "API Modules",
        "Validation and Security",
        "Screenshots and Results",
        "Conclusion",
    ]),
    ("Problem Statement", [
        "Manual internship tracking lacks transparency.",
        "Approval and review cycles are delayed.",
        "Coordination between student, faculty, and admin is inefficient.",
    ]),
    ("Objectives", [
        "Digitize internship lifecycle end-to-end.",
        "Implement secure role-based access control.",
        "Enable report submissions with mentor feedback.",
        "Provide analytics and audit-ready records.",
    ]),
    ("Technology Stack", [
        "Frontend: React, Vite, Tailwind CSS, Framer Motion",
        "Backend: Node.js, Express",
        "Database: MongoDB, Mongoose",
        "Auth: JWT, bcryptjs",
        "File Upload: Multer",
        "Testing: Jest, Supertest",
    ]),
    ("System Architecture", [
        "3-tier architecture: Frontend -> API -> Database",
        "JWT-based auth for protected endpoints",
        "Validation and role checks at middleware layer",
        "Audit logs for traceability",
        "Diagram source: docs/diagrams/architecture.mmd",
    ]),
    ("DFD Level 0", [
        "External entities: Student, Faculty, Admin",
        "Core process: SITS System",
        "Data stores: Users, Internships, Reports, Audit",
        "Diagram source: docs/diagrams/dfd_level0.mmd",
    ]),
    ("DFD Level 1", [
        "Authentication, Internship, Reporting, Admin, Notification processes",
        "Shows flow between entities, processes, and stores",
        "Diagram source: docs/diagrams/dfd_level1.mmd",
    ]),
    ("ER Diagram", [
        "Entities: User, StudentProfile, Internship, Report, Feedback, AuditLog",
        "Relationships mapped using ObjectId references",
        "Diagram source: docs/diagrams/er_diagram.mmd",
    ]),
    ("Database Schema - User & StudentProfile", [
        "User: name, email(unique), password(hash), role, isActive",
        "StudentProfile: one-to-one with User",
        "Stores academic and portfolio fields",
    ]),
    ("Database Schema - Internship & Report", [
        "Internship: student, company, role, dates, status, mentor",
        "Report: internship, weekNumber, content, attachment, status",
        "Review metadata: feedback, reviewedBy, reviewedAt",
    ]),
    ("Database Schema - Feedback & AuditLog", [
        "Feedback links report, internship, student, faculty",
        "AuditLog tracks actor/action/entity for accountability",
        "Supports review and compliance use-cases",
    ]),
    ("API Module Overview", [
        "Auth: /api/auth/register, /api/auth/login",
        "Internship: create, list, approve/reject, certificate upload",
        "Reports: submit, fetch by internship, feedback update",
        "Admin: users CRUD, assign mentor, stats",
        "Profile: get/update",
    ]),
    ("Validation and Security", [
        "Backend validation via express-validator",
        "Frontend inline field validation",
        "JWT authentication and role authorization",
        "ObjectId validation on route parameters",
        "Password hashing with bcryptjs",
    ]),
    ("Notification Design", [
        "Polling interval: 45 seconds",
        "Student: status + feedback alerts",
        "Faculty: assigned students only",
        "Admin: global pending overview",
        "Scoped persistence in localStorage per user",
    ]),
    ("UI Screenshots - Student Flow", [
        "Insert: login, register, student dashboard",
        "Insert: add internship and reports page",
        "Place screenshot files from docs/screenshots",
    ]),
    ("UI Screenshots - Faculty/Admin Flow", [
        "Insert: faculty dashboard, admin dashboard",
        "Insert: mentor assignment and notifications panel",
        "Use clear annotations where needed",
    ]),
    ("Testing and Build Status", [
        "Integration tests validate RBAC and key workflows",
        "Frontend build compiled successfully",
        "Validated forms reduce invalid submissions",
    ]),
    ("Results, Limitations, Future Scope", [
        "Result: transparent and structured internship process",
        "Limitation: local file storage, polling-based alerts",
        "Future: cloud storage, push notifications, advanced analytics",
    ]),
    ("Conclusion", [
        "SITS provides a secure and scalable internship management foundation.",
        "Project demonstrates full-stack engineering, schema design, and workflow automation.",
        "Ready for institutional pilot with minor enhancements.",
    ]),
    ("Thank You", [
        "Questions and Discussion",
    ]),
]

prs = Presentation()

for title, bullets in slides_data:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22 if idx == 0 else 18)

# Add a final slide with screenshot checklist
layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(layout)
slide.shapes.title.text = "Screenshot Checklist"

tf = slide.shapes.placeholders[1].text_frame
tf.clear()
checklist = [
    "ss01_login.png",
    "ss02_register.png",
    "ss03_student_dashboard.png",
    "ss04_add_internship.png",
    "ss05_reports_page.png",
    "ss06_faculty_dashboard.png",
    "ss07_admin_dashboard.png",
    "ss08_assign_mentor.png",
    "ss09_profile_page.png",
    "ss10_notifications_panel.png",
]

for i, item in enumerate(checklist):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.level = 0
    p.font.size = Pt(18)

OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT_PATH))
print(f"PPT generated: {OUT_PATH}")
